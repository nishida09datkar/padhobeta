from pptx import Presentation
from models.schemas import ParsedDocument


def parse_pptx(file_path: str) -> ParsedDocument:
    prs = Presentation(file_path)

    title = None
    if prs.core_properties.title:
        title = prs.core_properties.title

    lines = []
    slide_count = len(prs.slides)

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    level = para.level if para.level else 0
                    if level == 0:
                        lines.append(f"## {text}")
                    else:
                        indent = "  " * (level - 1)
                        lines.append(f"{indent}- {text}")

            if shape.has_table:
                table = shape.table
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table.rows[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")

        lines.append("")

    content = "\n".join(lines)

    if not content.strip():
        content = "[PPTX file appears to be empty or contains no readable text.]"

    return ParsedDocument(
        content=content,
        title=title,
        doc_type="pptx",
        page_count=slide_count,
        metadata={
            "slide_count": slide_count,
        },
    )
